from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette couleurs ──────────────────────────────────────────
GREEN       = RGBColor(0x0A, 0x7C, 0x4E)   # vert principal
GREEN_LIGHT = RGBColor(0xD1, 0xFA, 0xE5)   # vert clair fond
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x1E, 0x29, 0x3B)   # quasi-noir
GRAY        = RGBColor(0x6B, 0x72, 0x80)
ACCENT      = RGBColor(0x63, 0x66, 0xF1)   # indigo pour accent

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # layout vide

# ─────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_multiline(slide, lines, l, t, w, h, size=14, color=DARK, bold_first=False, spacing=1.2):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
    return txb

# ─────────────────────────────────────────────────────────────
# DIAPO 1 — Titre
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)

# Fond vert foncé pleine page
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GREEN)

# Bande blanche décorative gauche
add_rect(slide, 0, 0, 0.55, 7.5, fill_rgb=WHITE)

# Pastille blanche en haut à droite
add_rect(slide, 10.5, 0.3, 2.5, 2.5, fill_rgb=RGBColor(0x0C, 0x96, 0x5F))

# Titre principal
add_text(slide, "Hawa", 1.2, 1.6, 9, 1.5, size=72, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Sous-titre
add_text(slide, "Chatbot de Santé Communautaire — Niger",
         1.2, 3.1, 10, 0.8, size=26, bold=False, color=RGBColor(0xA7, 0xF3, 0xD0), align=PP_ALIGN.LEFT)

# Description
add_text(slide, "Assistant IA sur WhatsApp pour la vaccination et la santé du nourrisson",
         1.2, 3.9, 10, 0.7, size=18, color=RGBColor(0xD1, 0xFA, 0xE5), align=PP_ALIGN.LEFT)

# Ligne décorative
add_rect(slide, 1.2, 3.0, 6, 0.04, fill_rgb=RGBColor(0x34, 0xD3, 0x99))

# Tags en bas
for i, tag in enumerate(["Ministère de la Santé · Niger", "Programme PEV", "IA Générative"]):
    add_rect(slide, 1.2 + i * 3.5, 5.8, 3.1, 0.45, fill_rgb=RGBColor(0x05, 0x6B, 0x40))
    add_text(slide, tag, 1.25 + i * 3.5, 5.83, 3.0, 0.4, size=12, color=RGBColor(0xA7, 0xF3, 0xD0), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────
# DIAPO 2 — Contexte & Problématique
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, fill_rgb=GREEN)
add_rect(slide, 0, 1.2, 0.12, 6.3, fill_rgb=GREEN)

add_text(slide, "Contexte & Problématique", 0.4, 0.25, 10, 0.75, size=28, bold=True, color=WHITE)
add_text(slide, "02", 11.8, 0.15, 1.2, 0.8, size=36, bold=True, color=RGBColor(0x34, 0xD3, 0x99))

# Carte gauche
add_rect(slide, 0.4, 1.5, 5.8, 2.5, fill_rgb=GREEN_LIGHT)
add_text(slide, "🌍  Situation au Niger", 0.6, 1.6, 5.4, 0.55, size=14, bold=True, color=GREEN)
add_multiline(slide, [
    "• Taux de mortalité infantile élevé",
    "• Couverture vaccinale insuffisante en zones rurales",
    "• Barrière linguistique : majorité Hausa-phone",
    "• Accès limité aux agents de santé communautaires",
], 0.6, 2.15, 5.5, 1.7, size=13, color=DARK)

# Carte droite
add_rect(slide, 6.6, 1.5, 6.3, 2.5, fill_rgb=RGBColor(0xEE, 0xF2, 0xFF))
add_text(slide, "📱  Opportunité numérique", 6.8, 1.6, 5.8, 0.55, size=14, bold=True, color=ACCENT)
add_multiline(slide, [
    "• WhatsApp : messagerie la plus utilisée au Niger",
    "• Pénétration mobile en forte croissance",
    "• Les familles envoient des messages vocaux",
    "• Besoin d'information santé fiable et accessible",
], 6.8, 2.15, 6.0, 1.7, size=13, color=DARK)

# Bandeau bas
add_rect(slide, 0.4, 4.3, 12.5, 1.2, fill_rgb=RGBColor(0xF0, 0xFD, 0xF4))
add_rect(slide, 0.4, 4.3, 0.08, 1.2, fill_rgb=GREEN)
add_text(slide, "💡  Solution : un chatbot IA bilingue (Français / Hausa) accessible sur WhatsApp, disponible 24h/24, qui répond aux questions de vaccination et de santé du nourrisson.",
         0.65, 4.4, 12.0, 1.0, size=14, bold=False, color=DARK)

# ─────────────────────────────────────────────────────────────
# DIAPO 3 — Objectifs du projet
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, fill_rgb=GREEN)
add_rect(slide, 0, 1.2, 0.12, 6.3, fill_rgb=GREEN)

add_text(slide, "Objectifs du Projet", 0.4, 0.25, 10, 0.75, size=28, bold=True, color=WHITE)
add_text(slide, "03", 11.8, 0.15, 1.2, 0.8, size=36, bold=True, color=RGBColor(0x34, 0xD3, 0x99))

objectives = [
    ("🎯", "Informer sur la vaccination",
     "Répondre aux questions sur le calendrier PEV Niger, les effets secondaires et les vaccins disponibles gratuitement."),
    ("👶", "Santé du nourrisson",
     "Accompagner les parents sur la nutrition, les maladies courantes et les signes d'alerte chez l'enfant."),
    ("🗣️", "Accessibilité bilingue",
     "Comprendre et répondre en Français et en Hausa, avec détection automatique de la langue."),
    ("🎙️", "Messages vocaux",
     "Transcrire les messages audio des utilisateurs et répondre par synthèse vocale (TTS)."),
    ("👨‍⚕️", "Tableau de bord de gestion",
     "Permettre aux agents de santé de superviser les conversations et prendre la main si nécessaire."),
    ("📢", "Diffusions ciblées",
     "Envoyer des messages de prévention groupés à des listes de contacts WhatsApp."),
]

cols = [(0.4, 1.4), (4.6, 1.4), (8.8, 1.4),
        (0.4, 3.8), (4.6, 3.8), (8.8, 3.8)]

for i, (icon, title, desc) in enumerate(objectives):
    x, y = cols[i]
    add_rect(slide, x, y, 3.8, 2.1, fill_rgb=GREEN_LIGHT)
    add_rect(slide, x, y, 0.07, 2.1, fill_rgb=GREEN)
    add_text(slide, f"{icon}  {title}", x + 0.2, y + 0.12, 3.5, 0.5, size=13, bold=True, color=GREEN)
    add_text(slide, desc, x + 0.2, y + 0.6, 3.5, 1.35, size=11, color=DARK)

# ─────────────────────────────────────────────────────────────
# DIAPO 4 — Architecture technique
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, fill_rgb=GREEN)
add_rect(slide, 0, 1.2, 0.12, 6.3, fill_rgb=GREEN)

add_text(slide, "Architecture Technique", 0.4, 0.25, 10, 0.75, size=28, bold=True, color=WHITE)
add_text(slide, "04", 11.8, 0.15, 1.2, 0.8, size=36, bold=True, color=RGBColor(0x34, 0xD3, 0x99))

# Flèche centrale horizontale
add_rect(slide, 1.0, 3.5, 11.5, 0.06, fill_rgb=RGBColor(0xD1, 0xD5, 0xDB))

layers = [
    (0.4,  1.35, "📱\nWhatsApp\nUtilisateur",     RGBColor(0x25, 0xD3, 0x66), WHITE),
    (3.2,  1.35, "⚙️\nServeur\nNode.js / Express", GREEN,                       WHITE),
    (6.0,  1.35, "🤖\nOpenAI\nGPT-4o",             ACCENT,                      WHITE),
    (8.8,  1.35, "🗄️\nMongoDB\nBase de données",   RGBColor(0x10, 0xB9, 0x81),  WHITE),
    (11.4, 1.35, "🖥️\nDashboard\nReact",           RGBColor(0xF5, 0x9E, 0x0B),  WHITE),
]

for x, y, label, bg, fg in layers:
    add_rect(slide, x, y, 2.2, 1.8, fill_rgb=bg)
    add_text(slide, label, x + 0.1, y + 0.2, 2.0, 1.4, size=13, bold=True, color=fg, align=PP_ALIGN.CENTER)

# Flèches entre les blocs
for xi in [2.62, 5.42, 8.22, 11.02]:
    add_text(slide, "→", xi, 2.0, 0.4, 0.45, size=18, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Description des flux
flows = [
    (0.4,  3.75, "Message texte\nou vocal"),
    (3.2,  3.75, "Webhook Meta\nTranscription STT"),
    (6.0,  3.75, "Génération de\nla réponse IA"),
    (8.8,  3.75, "Stockage\nconversations"),
    (11.4, 3.75, "Supervision\nadmin"),
]
for x, y, txt in flows:
    add_text(slide, txt, x, y, 2.2, 0.7, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# Services tiers
add_text(slide, "Services tiers", 0.4, 4.65, 5, 0.4, size=13, bold=True, color=DARK)
services = [
    (0.4,  5.15, "☁️  Cloudinary",     "Stockage audio MP3",       RGBColor(0xFE, 0xF3, 0xC7)),
    (3.5,  5.15, "🎤  Google STT",     "Transcription vocale",     RGBColor(0xEE, 0xF2, 0xFF)),
    (6.6,  5.15, "🔊  ElevenLabs/OAI", "Synthèse vocale TTS",      RGBColor(0xFD, 0xF2, 0xF8)),
    (9.7,  5.15, "📊  Meta API",       "Envoi messages WhatsApp",  GREEN_LIGHT),
]
for x, y, title, desc, bg in services:
    add_rect(slide, x, y, 3.0, 1.2, fill_rgb=bg)
    add_text(slide, title, x + 0.15, y + 0.12, 2.7, 0.45, size=12, bold=True, color=DARK)
    add_text(slide, desc,  x + 0.15, y + 0.58, 2.7, 0.45, size=11, color=GRAY)

# ─────────────────────────────────────────────────────────────
# DIAPO 5 — Stack technique
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, fill_rgb=GREEN)
add_rect(slide, 0, 1.2, 0.12, 6.3, fill_rgb=GREEN)

add_text(slide, "Stack Technologique", 0.4, 0.25, 10, 0.75, size=28, bold=True, color=WHITE)
add_text(slide, "05", 11.8, 0.15, 1.2, 0.8, size=36, bold=True, color=RGBColor(0x34, 0xD3, 0x99))

categories = [
    ("Backend", GREEN, [
        "Node.js  —  runtime JavaScript",
        "Express.js  —  serveur HTTP & API REST",
        "Mongoose  —  ODM MongoDB",
        "JWT  —  authentification sécurisée",
        "Multer  —  upload de fichiers",
        "node-cron  —  tâches planifiées",
    ]),
    ("Intelligence Artificielle", ACCENT, [
        "OpenAI GPT-4o  —  génération des réponses",
        "Whisper (OpenAI)  —  transcription audio (STT)",
        "Google Cloud Speech  —  STT alternatif",
        "ElevenLabs / OpenAI TTS  —  synthèse vocale",
        "franc  —  détection de langue statistique",
    ]),
    ("Frontend (Dashboard)", RGBColor(0xF5, 0x9E, 0x0B), [
        "React 18  —  interface utilisateur",
        "Vite  —  bundler ultra-rapide",
        "React Router v6  —  navigation SPA",
        "Bootstrap Icons  —  icônes UI",
    ]),
    ("Infrastructure & Cloud", RGBColor(0x10, 0xB9, 0x81), [
        "MongoDB Atlas  —  base de données cloud",
        "Cloudinary  —  stockage médias audio",
        "Meta (WhatsApp) API  —  messagerie",
        "dotenv  —  gestion des secrets",
    ]),
]

positions = [(0.4, 1.35), (3.55, 1.35), (6.7, 1.35), (9.85, 1.35)]

for (x, y), (cat, color, items) in zip(positions, categories):
    add_rect(slide, x, y, 3.0, 5.8, fill_rgb=RGBColor(0xF9, 0xFA, 0xFB))
    add_rect(slide, x, y, 3.0, 0.55, fill_rgb=color)
    add_text(slide, cat, x + 0.15, y + 0.08, 2.7, 0.42, size=13, bold=True, color=WHITE)
    for j, item in enumerate(items):
        add_text(slide, f"• {item}", x + 0.15, y + 0.7 + j * 0.76, 2.75, 0.65, size=11, color=DARK)

# ─────────────────────────────────────────────────────────────
# DIAPO 6 — Fonctionnalités clés
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, fill_rgb=GREEN)
add_rect(slide, 0, 1.2, 0.12, 6.3, fill_rgb=GREEN)

add_text(slide, "Fonctionnalités Clés", 0.4, 0.25, 10, 0.75, size=28, bold=True, color=WHITE)
add_text(slide, "06", 11.8, 0.15, 1.2, 0.8, size=36, bold=True, color=RGBColor(0x34, 0xD3, 0x99))

features = [
    ("🎙️ Messages vocaux",
     "L'utilisateur envoie un vocal WhatsApp → transcription automatique via Whisper → réponse IA → lecture TTS en retour",
     GREEN_LIGHT, GREEN),
    ("🌐 Bilingue Français / Hausa",
     "Détection automatique de la langue (franc + dictionnaire) · Réponse dans la langue du message · Support des caractères Hausa spéciaux (ƙ, ɗ, ɓ)",
     RGBColor(0xEE, 0xF2, 0xFF), ACCENT),
    ("🛡️ Mode Humain",
     "L'agent peut basculer en \"Mode Humain\" pour reprendre la conversation directement depuis le dashboard · Retour automatique au mode IA possible",
     RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0xD9, 0x77, 0x06)),
    ("📢 Diffusions groupées",
     "Envoi de messages de prévention vaccinale à des listes de contacts · Support texte, image, audio, vidéo et documents PDF",
     RGBColor(0xFD, 0xF2, 0xF8), RGBColor(0xEC, 0x48, 0x99)),
    ("📍 Géolocalisation",
     "L'utilisateur peut partager sa position pour trouver le centre de santé le plus proche (CSI)",
     GREEN_LIGHT, GREEN),
    ("🔒 Gestion des accès",
     "Système de rôles (Admin / Staff) · Authentification JWT · Blocage de contacts · Gestion des utilisateurs du dashboard",
     RGBColor(0xEE, 0xF2, 0xFF), ACCENT),
]

for i, (title, desc, bg, accent_col) in enumerate(features):
    row, col = divmod(i, 2)
    x = 0.4 + col * 6.4
    y = 1.45 + row * 1.8
    add_rect(slide, x, y, 6.1, 1.6, fill_rgb=bg)
    add_rect(slide, x, y, 0.07, 1.6, fill_rgb=accent_col)
    add_text(slide, title, x + 0.2, y + 0.1, 5.7, 0.5, size=13, bold=True, color=accent_col)
    add_text(slide, desc,  x + 0.2, y + 0.6, 5.7, 0.9, size=11, color=DARK)

# ─────────────────────────────────────────────────────────────
# DIAPO 7 — Tableau de bord
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, fill_rgb=GREEN)
add_rect(slide, 0, 1.2, 0.12, 6.3, fill_rgb=GREEN)

add_text(slide, "Tableau de Bord Administratif", 0.4, 0.25, 10, 0.75, size=28, bold=True, color=WHITE)
add_text(slide, "07", 11.8, 0.15, 1.2, 0.8, size=36, bold=True, color=RGBColor(0x34, 0xD3, 0x99))

# Mockup dashboard
add_rect(slide, 0.4, 1.35, 12.5, 5.8, fill_rgb=RGBColor(0xF3, 0xF4, 0xF6))

# Sidebar
add_rect(slide, 0.4, 1.35, 2.0, 5.8, fill_rgb=DARK)
menu_items = ["Dashboard", "Discussions", "Contacts", "Diffusions", "Métadonnées", "Paramètres"]
for i, item in enumerate(menu_items):
    if item == "Discussions":
        add_rect(slide, 0.4, 1.35 + 0.65 + i * 0.72, 2.0, 0.52, fill_rgb=GREEN)
    add_text(slide, item, 0.55, 1.5 + i * 0.72, 1.7, 0.45, size=11, color=WHITE, bold=(item == "Discussions"))

# Zone principale
add_rect(slide, 2.4, 1.35, 10.5, 0.55, fill_rgb=WHITE)
add_text(slide, "Discussions — Historique des conversations WhatsApp", 2.55, 1.42, 9, 0.4, size=12, bold=True, color=DARK)

# Stats cards
stats = [("💬 Msgs reçus", "1 247"), ("🎙️ Audio", "389"), ("📝 Texte", "858"), ("👥 Contacts", "214")]
for i, (label, val) in enumerate(stats):
    add_rect(slide, 2.5 + i * 2.55, 2.0, 2.35, 1.1, fill_rgb=WHITE)
    add_text(slide, val, 2.65 + i * 2.55, 2.05, 2.0, 0.55, size=22, bold=True, color=GREEN)
    add_text(slide, label, 2.65 + i * 2.55, 2.6, 2.0, 0.35, size=10, color=GRAY)

# Table discussions
add_rect(slide, 2.4, 3.25, 10.5, 0.45, fill_rgb=RGBColor(0xE5, 0xE7, 0xEB))
headers = ["Contact", "Numéro", "Statut", "Langue", "Actions"]
widths  = [2.0, 2.0, 1.8, 1.5, 2.0]
xi = 2.5
for h, w in zip(headers, widths):
    add_text(slide, h, xi, 3.28, w, 0.38, size=10, bold=True, color=DARK)
    xi += w + 0.05

rows_data = [
    ("Fatima Moussa", "+227 96 XX XX XX", "🟢 Ouvert",       "Hausa",   "Ouvrir  + Contact"),
    ("Aissatou B.",   "+227 93 XX XX XX", "🔵 Mode Humain",  "Français","Ouvrir  ✓ Enregistré"),
    ("Hadiza Idrissa","+227 91 XX XX XX", "🔴 Fermé",        "Hausa",   "Ouvrir  + Contact"),
]
for ri, row in enumerate(rows_data):
    bg = WHITE if ri % 2 == 0 else RGBColor(0xF9, 0xFA, 0xFB)
    add_rect(slide, 2.4, 3.72 + ri * 0.6, 10.5, 0.55, fill_rgb=bg)
    xi = 2.5
    for val, w in zip(row, widths):
        add_text(slide, val, xi, 3.76 + ri * 0.6, w, 0.42, size=10, color=DARK)
        xi += w + 0.05

# ─────────────────────────────────────────────────────────────
# DIAPO 8 — Flux conversation WhatsApp
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, fill_rgb=GREEN)
add_rect(slide, 0, 1.2, 0.12, 6.3, fill_rgb=GREEN)

add_text(slide, "Flux d'une Conversation WhatsApp", 0.4, 0.25, 10, 0.75, size=28, bold=True, color=WHITE)
add_text(slide, "08", 11.8, 0.15, 1.2, 0.8, size=36, bold=True, color=RGBColor(0x34, 0xD3, 0x99))

steps = [
    ("1", "Message\narrive",     "Texte ou vocal\nvia WhatsApp",  GREEN),
    ("2", "Détection\nlangue",   "Français ou\nHausa",           ACCENT),
    ("3", "STT\n(si vocal)",     "Transcription\nWhisper/Google", RGBColor(0xF5, 0x9E, 0x0B)),
    ("4", "Contexte\nIA",        "Historique +\nPrompt système",  RGBColor(0x10, 0xB9, 0x81)),
    ("5", "GPT-4o\nRépond",      "Réponse courte\n3-5 phrases",   RGBColor(0xEC, 0x48, 0x99)),
    ("6", "TTS\n(si vocal)",     "Synthèse\naudio MP3",           RGBColor(0x6B, 0x72, 0x80)),
    ("7", "Envoi\nWhatsApp",     "Réponse au\ncontact",           GREEN),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.5 + i * 1.82
    # Connecteur
    if i < len(steps) - 1:
        add_rect(slide, x + 1.35, 2.55, 0.47, 0.06, fill_rgb=RGBColor(0xD1, 0xD5, 0xDB))
    # Cercle numéroté
    add_rect(slide, x, 1.4, 1.3, 1.3, fill_rgb=color)
    add_text(slide, num, x, 1.4, 1.3, 1.3, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x - 0.05, 2.8, 1.4, 0.7, size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(slide, desc,  x - 0.05, 3.5, 1.4, 0.7, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# Cas particulier stockage
add_rect(slide, 0.5, 4.5, 12.3, 1.5, fill_rgb=GREEN_LIGHT)
add_rect(slide, 0.5, 4.5, 0.08, 1.5, fill_rgb=GREEN)
add_text(slide, "💾  Stockage systématique", 0.75, 4.6, 5, 0.45, size=13, bold=True, color=GREEN)
add_multiline(slide, [
    "• Chaque message (entrant et sortant) est enregistré en MongoDB avec : type de contenu, langue, transcription, URL audio, métadonnées IA",
    "• Le contact est créé automatiquement à son premier message · Le contexte des 20 derniers échanges est passé à GPT pour une réponse cohérente",
], 0.75, 5.05, 12.0, 0.85, size=11, color=DARK)

# ─────────────────────────────────────────────────────────────
# DIAPO 9 — Impact & Bilan
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
add_rect(slide, 0, 0, 13.33, 1.2, fill_rgb=GREEN)
add_rect(slide, 0, 1.2, 0.12, 6.3, fill_rgb=GREEN)

add_text(slide, "Impact & Bilan", 0.4, 0.25, 10, 0.75, size=28, bold=True, color=WHITE)
add_text(slide, "09", 11.8, 0.15, 1.2, 0.8, size=36, bold=True, color=RGBColor(0x34, 0xD3, 0x99))

impacts = [
    ("✅", "Accessibilité maximale",   "Aucune application à installer · WhatsApp déjà présent sur le téléphone des familles"),
    ("✅", "Barrière linguistique levée", "Support natif du Hausa, langue maternelle de plus de 50 % de la population nigérienne"),
    ("✅", "Disponibilité 24h/24",      "Réponse immédiate à toute heure, même en l'absence d'agent de santé"),
    ("✅", "Informations validées",      "Contenu basé sur le calendrier PEV officiel du Ministère de la Santé du Niger"),
    ("✅", "Supervision humaine",        "Les agents de santé gardent le contrôle total via le tableau de bord"),
    ("✅", "Scalabilité",               "Un seul serveur peut gérer des milliers de conversations simultanément"),
]

for i, (icon, title, desc) in enumerate(impacts):
    row, col = divmod(i, 2)
    x = 0.4 + col * 6.4
    y = 1.45 + row * 1.7
    add_rect(slide, x, y, 6.1, 1.5, fill_rgb=GREEN_LIGHT if col == 0 else RGBColor(0xF0, 0xFD, 0xF4))
    add_text(slide, f"{icon}  {title}", x + 0.2, y + 0.1, 5.7, 0.5, size=13, bold=True, color=GREEN)
    add_text(slide, desc, x + 0.2, y + 0.65, 5.7, 0.75, size=11, color=DARK)

# ─────────────────────────────────────────────────────────────
# DIAPO 10 — Merci / Contact
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GREEN)
add_rect(slide, 0, 0, 0.55, 7.5, fill_rgb=WHITE)
add_rect(slide, 0, 6.8, 13.33, 0.7, fill_rgb=RGBColor(0x05, 0x6B, 0x40))

add_text(slide, "Merci", 1.2, 1.5, 10, 1.4, size=80, bold=True, color=WHITE)
add_text(slide, "Hawa — Chatbot IA de Santé Communautaire",
         1.2, 3.1, 10, 0.7, size=22, color=RGBColor(0xA7, 0xF3, 0xD0))
add_rect(slide, 1.2, 3.0, 6, 0.04, fill_rgb=RGBColor(0x34, 0xD3, 0x99))

add_multiline(slide, [
    "🌐  Projet e-Santé Niger",
    "📱  WhatsApp · Node.js · React · OpenAI · MongoDB",
    "🤝  En partenariat avec le Ministère de la Santé Publique du Niger",
], 1.2, 3.85, 11, 1.8, size=15, color=RGBColor(0xD1, 0xFA, 0xE5))

add_text(slide, "Des questions ?", 1.2, 5.8, 10, 0.6, size=16, color=RGBColor(0x6E, 0xE7, 0xB7))

# ── Sauvegarde ────────────────────────────────────────────────
OUTPUT = r"C:\Users\ABOU\Desktop\e-sante\project-node-chatbot\Hawa_Presentation.pptx"
prs.save(OUTPUT)
print(f"Fichier généré : {OUTPUT}")
